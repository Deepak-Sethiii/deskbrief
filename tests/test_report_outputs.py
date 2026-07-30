"""Stage 6 coverage: chart rendering and the one-slide deck."""

from __future__ import annotations

import datetime as dt

import pandas as pd
import pytest
from pptx import Presentation
from pptx.util import Inches

from src.report.charts import normalised_price_chart
from src.report.deck import TABLE_COLUMNS, build_deck


@pytest.fixture()
def prices():
    days = pd.bdate_range("2026-01-01", periods=60)
    rows = []
    for ticker, start, drift in (("A.NS", 100.0, 0.5), ("B.NS", 2000.0, -3.0)):
        for index, day in enumerate(days):
            rows.append({"ticker": ticker, "date": day, "close": start + drift * index})
    return pd.DataFrame(rows)


@pytest.fixture()
def comps():
    return pd.DataFrame(
        {
            "sector": ["Technology", "Energy"],
            "last_close": [2446.6, 1275.9],
            "ret_1d": [0.0203, -0.0065],
            "ret_5d": [0.1079, -0.0099],
            "ret_21d": [float("nan"), -0.0139],
            "realised_vol_30d": [0.3536, 0.175],
            "pct_from_52w_high": [-0.2365, -0.0195],
            "pe": [17.65, float("nan")],
        },
        index=pd.Index(["TCS.NS", "RELIANCE.NS"], name="ticker"),
    )


@pytest.fixture()
def commentary():
    return {
        "market_tone": "Breadth was risk-on.",
        "bullets": [f"Observation {i}" for i in range(8)],  # more than MAX_BULLETS
        "watch_items": ["Watch one", "Watch two", "Watch three", "Watch four"],
        "source": "unit-test",
    }


# --------------------------------------------------------------------------
# charts
# --------------------------------------------------------------------------

def test_chart_is_written(prices, tmp_path):
    out = normalised_price_chart(prices, output_path=tmp_path / "c.png")
    assert out is not None and out.exists() and out.stat().st_size > 5_000


def test_compact_variant_has_a_different_aspect(prices, tmp_path):
    from PIL import Image

    wide = normalised_price_chart(prices, output_path=tmp_path / "w.png")
    compact = normalised_price_chart(prices, output_path=tmp_path / "s.png", compact=True)
    wide_ratio = Image.open(wide).width / Image.open(wide).height
    compact_ratio = Image.open(compact).width / Image.open(compact).height
    assert compact_ratio > wide_ratio  # short and wide, for the slide


def test_chart_returns_none_rather_than_raising_on_no_data(tmp_path):
    assert normalised_price_chart(pd.DataFrame(), output_path=tmp_path / "x.png") is None
    empty = pd.DataFrame(columns=["ticker", "date", "close"])
    assert normalised_price_chart(empty, output_path=tmp_path / "y.png") is None


def test_chart_survives_a_ticker_with_all_nan_closes(prices, tmp_path):
    broken = pd.concat(
        [prices, pd.DataFrame({"ticker": ["C.NS"] * 3,
                               "date": pd.bdate_range("2026-01-01", periods=3),
                               "close": [float("nan")] * 3})]
    )
    assert normalised_price_chart(broken, output_path=tmp_path / "z.png") is not None


# --------------------------------------------------------------------------
# deck
# --------------------------------------------------------------------------

def test_deck_has_exactly_one_slide(comps, commentary, tmp_path):
    out = build_deck(comps, commentary, None, output_dir=tmp_path)
    assert len(Presentation(out).slides) == 1


def test_deck_filename_matches_the_workbook_convention(comps, tmp_path):
    out = build_deck(comps, None, None, output_dir=tmp_path,
                     generated_at=dt.datetime(2026, 7, 31, 14, 5))
    assert out.name == "deskbrief_20260731_1405.pptx"


def test_deck_is_widescreen(comps, tmp_path):
    presentation = Presentation(build_deck(comps, None, None, output_dir=tmp_path))
    assert presentation.slide_width == Inches(13.333)
    assert presentation.slide_height == Inches(7.5)


def test_deck_uses_the_blank_layout_so_no_placeholder_junk(comps, tmp_path):
    slide = Presentation(build_deck(comps, None, None, output_dir=tmp_path)).slides[0]
    texts = [
        shape.text_frame.text
        for shape in slide.shapes
        if shape.has_text_frame
    ]
    joined = " ".join(texts)
    assert "Click to add" not in joined


def test_deck_table_shape_and_contents(comps, commentary, tmp_path):
    slide = Presentation(build_deck(comps, commentary, None, output_dir=tmp_path)).slides[0]
    tables = [s.table for s in slide.shapes if s.has_table]
    assert len(tables) == 1

    table = tables[0]
    assert len(table.rows) == len(comps) + 1
    assert len(table.columns) == len(TABLE_COLUMNS)
    assert table.cell(0, 0).text == "Ticker"
    assert table.cell(1, 0).text == "TCS"  # .NS suffix stripped for the slide
    assert table.cell(1, 2).text == "+2.03%"
    assert table.cell(1, 5).text == "-23.6%"  # -23.65 rounds half-to-even
    assert table.cell(1, 4).text == "-"  # NaN 21d return renders as a dash


def test_deck_caps_the_bullet_count(comps, commentary, tmp_path):
    from src.report.deck import MAX_BULLETS

    slide = Presentation(build_deck(comps, commentary, None, output_dir=tmp_path)).slides[0]
    joined = " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
    assert f"Observation {MAX_BULLETS - 1}" in joined
    assert f"Observation {MAX_BULLETS}" not in joined


def test_deck_embeds_the_chart_when_present(comps, commentary, prices, tmp_path):
    chart = normalised_price_chart(prices, output_path=tmp_path / "c.png", compact=True)
    slide = Presentation(
        build_deck(comps, commentary, chart, output_dir=tmp_path)
    ).slides[0]
    assert any(shape.shape_type == 13 for shape in slide.shapes)  # 13 = PICTURE


def test_deck_says_so_when_the_chart_is_missing(comps, tmp_path):
    slide = Presentation(
        build_deck(comps, None, tmp_path / "does_not_exist.png", output_dir=tmp_path)
    ).slides[0]
    joined = " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
    assert "Chart unavailable" in joined


def test_deck_footer_names_the_commentary_source(comps, commentary, tmp_path):
    slide = Presentation(build_deck(comps, commentary, None, output_dir=tmp_path)).slides[0]
    joined = " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
    assert "unit-test" in joined
    assert "not verified entity links" in joined


def test_deck_builds_with_no_data_at_all(tmp_path):
    out = build_deck(pd.DataFrame(), None, None, output_dir=tmp_path)
    slide = Presentation(out).slides[0]
    joined = " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
    assert "No comparable data available." in joined
