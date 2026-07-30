"""PowerPoint tear-sheet: exactly one slide.

Built on the blank layout (index 6) rather than a title/content layout, because
the stock layouts drag in "Click to add title" placeholders that survive into
the saved file and look like an unfinished deck. Every box here is positioned
explicitly, so what you see is only what we put there.

One slide is a deliberate constraint: a tear-sheet is something an analyst
glances at, not a deck they page through.
"""

from __future__ import annotations

import datetime as dt
import logging
from pathlib import Path
from typing import Any, Mapping

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

log = logging.getLogger(__name__)

PROJECT_ROOT = Path(__file__).resolve().parents[2]
OUTPUT_DIR = PROJECT_ROOT / "output"

SLIDE_WIDTH = Inches(13.333)  # 16:9
SLIDE_HEIGHT = Inches(7.5)

INK = RGBColor(0x1F, 0x2A, 0x37)
MUTED = RGBColor(0x6B, 0x72, 0x80)
ACCENT = RGBColor(0x1F, 0x3A, 0x5B)
HEADER_TEXT = RGBColor(0xFF, 0xFF, 0xFF)
BAND = RGBColor(0xF3, 0xF5, 0xF7)
POSITIVE = RGBColor(0x0B, 0x7A, 0x3B)
NEGATIVE = RGBColor(0xB3, 0x26, 0x1E)

TABLE_COLUMNS: tuple[tuple[str, str, float], ...] = (
    # (header, comps column, width in inches)
    ("Ticker", "_index", 1.55),
    ("Last", "last_close", 1.00),
    ("1D", "ret_1d", 0.85),
    ("5D", "ret_5d", 0.85),
    ("21D", "ret_21d", 0.85),
    ("vs 52W", "pct_from_52w_high", 0.95),
    ("Vol30", "realised_vol_30d", 0.85),
    ("P/E", "pe", 0.80),
)
MAX_BULLETS = 5


def _format_cell(column: str, value: Any) -> tuple[str, RGBColor | None]:
    """Return the display string and an optional colour for signed columns."""
    if value is None or (isinstance(value, float) and pd.isna(value)) or pd.isna(value):
        return "-", None

    if column in ("ret_1d", "ret_5d", "ret_21d"):
        colour = POSITIVE if value > 0 else (NEGATIVE if value < 0 else None)
        return f"{value * 100:+.2f}%", colour
    if column == "pct_from_52w_high":
        return f"{value * 100:.1f}%", None
    if column == "realised_vol_30d":
        return f"{value * 100:.0f}%", None
    if column == "last_close":
        return f"{value:,.0f}", None
    if column == "pe":
        return f"{value:.1f}", None
    return str(value)


def _text_box(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Emu(0)
    frame.margin_top = frame.margin_bottom = Emu(0)
    return frame


def _style_run(run, *, size: int, bold: bool = False, colour: RGBColor = INK):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = colour
    run.font.name = "Calibri"


def build_deck(
    comps: pd.DataFrame,
    commentary: Mapping[str, Any] | None = None,
    chart_path: str | Path | None = None,
    *,
    output_dir: str | Path | None = None,
    generated_at: dt.datetime | None = None,
) -> Path:
    """Write output/deskbrief_YYYYMMDD_HHMM.pptx and return its path."""
    stamp = generated_at or dt.datetime.now()
    out_dir = Path(output_dir) if output_dir else OUTPUT_DIR
    out_dir.mkdir(parents=True, exist_ok=True)

    presentation = Presentation()
    presentation.slide_width = SLIDE_WIDTH
    presentation.slide_height = SLIDE_HEIGHT
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # 6 = blank

    _add_header(slide, comps, stamp)
    table_bottom = _add_table(slide, comps)
    _add_chart(slide, chart_path, top=table_bottom)
    _add_commentary(slide, commentary)
    _add_footer(slide, commentary, stamp)

    out_path = out_dir / f"deskbrief_{stamp:%Y%m%d_%H%M}.pptx"
    presentation.save(out_path)
    log.info("deck written: %s", out_path.resolve())
    return out_path.resolve()


def _add_header(slide, comps: pd.DataFrame, stamp: dt.datetime) -> None:
    frame = _text_box(slide, Inches(0.45), Inches(0.32), Inches(12.4), Inches(0.9))

    paragraph = frame.paragraphs[0]
    _style_run(paragraph.add_run(), size=26, bold=True, colour=ACCENT)
    paragraph.runs[0].text = "DeskBrief - Indian Large Cap Watchlist"

    subtitle = frame.add_paragraph()
    run = subtitle.add_run()
    run.text = (
        f"{len(comps)} names  |  generated {stamp:%d %b %Y %H:%M}  |  "
        "source: Yahoo Finance via yfinance, Indian financial RSS"
    )
    _style_run(run, size=11, colour=MUTED)

    line = slide.shapes.add_shape(1, Inches(0.45), Inches(1.24), Inches(12.44), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    line.shadow.inherit = False


def _add_table(slide, comps: pd.DataFrame):
    """Draw the comps table top-left. Returns the y coordinate just below it."""
    if comps is None or comps.empty:
        frame = _text_box(slide, Inches(0.45), Inches(1.6), Inches(7.7), Inches(0.4))
        _style_run(frame.paragraphs[0].add_run(), size=12, colour=MUTED)
        frame.paragraphs[0].runs[0].text = "No comparable data available."
        return Inches(2.2)

    rows = len(comps) + 1
    columns = len(TABLE_COLUMNS)
    total_width = Inches(sum(width for _, _, width in TABLE_COLUMNS))
    row_height = Inches(0.32)

    graphic_frame = slide.shapes.add_table(
        rows, columns, Inches(0.45), Inches(1.6), total_width, row_height * rows
    )
    table = graphic_frame.table
    # python-pptx applies a banded blue style by default; we overwrite every
    # cell's fill below so the table matches the rest of the slide.
    for index, (_, _, width) in enumerate(TABLE_COLUMNS):
        table.columns[index].width = Inches(width)

    for index, (header, _, _) in enumerate(TABLE_COLUMNS):
        cell = table.cell(0, index)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.LEFT if index == 0 else PP_ALIGN.RIGHT
        _style_run(paragraph.runs[0], size=10, bold=True, colour=HEADER_TEXT)

    for row_index, (ticker, record) in enumerate(comps.iterrows(), start=1):
        table.rows[row_index].height = row_height
        for column_index, (_, source, _) in enumerate(TABLE_COLUMNS):
            if source == "_index":
                text, colour = str(ticker).replace(".NS", ""), None
            else:
                text, colour = _format_cell(source, record.get(source))

            cell = table.cell(row_index, column_index)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = BAND if row_index % 2 else RGBColor(0xFF, 0xFF, 0xFF)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT if column_index == 0 else PP_ALIGN.RIGHT
            _style_run(
                paragraph.runs[0],
                size=10,
                bold=(column_index == 0),
                colour=colour or INK,
            )

    return Emu(graphic_frame.top + row_height * rows + Inches(0.3))


def _add_chart(slide, chart_path: str | Path | None, *, top) -> None:
    """Place the chart across the width of the left column, under the table."""
    left, width = Inches(0.45), Inches(7.7)
    if not chart_path or not Path(chart_path).exists():
        frame = _text_box(slide, left, top, width, Inches(0.4))
        _style_run(frame.paragraphs[0].add_run(), size=11, colour=MUTED)
        frame.paragraphs[0].runs[0].text = "Chart unavailable for this run."
        return

    slide.shapes.add_picture(str(chart_path), left, top, width=width)


def _add_commentary(slide, commentary: Mapping[str, Any] | None) -> None:
    left, top, width = Inches(8.45), Inches(1.6), Inches(4.45)
    frame = _text_box(slide, left, top, width, Inches(5.3))

    heading = frame.paragraphs[0]
    _style_run(heading.add_run(), size=13, bold=True, colour=ACCENT)
    heading.runs[0].text = "Commentary"

    if not commentary:
        paragraph = frame.add_paragraph()
        _style_run(paragraph.add_run(), size=10, colour=MUTED)
        paragraph.runs[0].text = "No commentary was generated for this run."
        return

    tone = commentary.get("market_tone")
    if tone:
        paragraph = frame.add_paragraph()
        paragraph.space_before = Pt(4)
        _style_run(paragraph.add_run(), size=10, bold=True, colour=INK)
        paragraph.runs[0].text = str(tone)

    for item in (commentary.get("bullets") or [])[:MAX_BULLETS]:
        paragraph = frame.add_paragraph()
        paragraph.space_before = Pt(5)
        run = paragraph.add_run()
        run.text = f"•  {item}"
        _style_run(run, size=10, colour=INK)

    watch = (commentary.get("watch_items") or [])[:3]
    if watch:
        paragraph = frame.add_paragraph()
        paragraph.space_before = Pt(10)
        _style_run(paragraph.add_run(), size=11, bold=True, colour=ACCENT)
        paragraph.runs[0].text = "Watch"
        for item in watch:
            line = frame.add_paragraph()
            line.space_before = Pt(4)
            run = line.add_run()
            run.text = f"•  {item}"
            _style_run(run, size=10, colour=MUTED)


def _add_footer(slide, commentary: Mapping[str, Any] | None, stamp: dt.datetime) -> None:
    frame = _text_box(slide, Inches(0.45), Inches(7.05), Inches(12.4), Inches(0.3))
    source = (commentary or {}).get("source", "n/a")
    run = frame.paragraphs[0].add_run()
    run.text = (
        f"Commentary source: {source}.  Headline ticker tags are substring matches, "
        "not verified entity links.  Not investment advice."
    )
    _style_run(run, size=8, colour=MUTED)
